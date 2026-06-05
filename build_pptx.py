# -*- coding: utf-8 -*-
"""Gera a apresentacao (pptx) do Trabalho 1 de VCPI — GTSRB."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

ACCENT = RGBColor(0x4C, 0x72, 0xB0)
DARK   = RGBColor(0x22, 0x2A, 0x35)
GREY   = RGBColor(0x55, 0x5B, 0x66)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xEE, 0xF1, 0xF6)

FIG = 'report/figs'
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bar(slide):
    """barra de cor no topo"""
    b = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    b.fill.solid(); b.fill.fore_color.rgb = ACCENT
    b.line.fill.background()
    return b


def txt(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def title_slide(slide, title, subtitle):
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
    band = slide.shapes.add_shape(1, 0, Inches(2.55), SW, Inches(0.08))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
    txt(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6), title,
        size=40, bold=True, color=WHITE)
    txt(slide, Inches(0.92), Inches(1.5), Inches(11.5), Inches(1.0), subtitle,
        size=20, color=RGBColor(0xBC, 0xCC, 0xE0))
    txt(slide, Inches(0.92), Inches(6.4), Inches(11.5), Inches(0.6),
        'Visão Computacional e Processamento de Imagem  ·  Trabalho Prático 1  ·  Junho 2026',
        size=14, color=RGBColor(0x90, 0x9C, 0xAD))


def content_slide(title, kicker=None):
    s = add_slide(); bar(s)
    if kicker:
        txt(s, Inches(0.6), Inches(0.32), Inches(12), Inches(0.4), kicker.upper(),
            size=12, bold=True, color=ACCENT)
        ty = Inches(0.62)
    else:
        ty = Inches(0.4)
    txt(s, Inches(0.6), ty, Inches(12.1), Inches(0.9), title, size=28, bold=True, color=DARK)
    return s


def bullets(slide, items, l=Inches(0.7), t=Inches(1.7), w=Inches(12), h=Inches(5.3),
            size=18, gap=10):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl = it
        else:
            text, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap)
        # marcador
        rb = p.add_run(); rb.text = ('•  ' if lvl == 0 else '–  ')
        rb.font.color.rgb = ACCENT if lvl == 0 else GREY
        rb.font.size = Pt(size); rb.font.bold = True
        # texto (suporta **negrito** simples por segmentos)
        for seg, b in _segments(text):
            r = p.add_run(); r.text = seg
            r.font.size = Pt(size - (2 if lvl else 0))
            r.font.bold = b
            r.font.color.rgb = DARK if not b else ACCENT
    return tb


def _segments(text):
    """divide por **...** em (texto, bold)"""
    out = []; parts = text.split('**')
    for i, p in enumerate(parts):
        if p:
            out.append((p, i % 2 == 1))
    return out


def image(slide, path, l, t, w=None, h=None):
    if w: return slide.shapes.add_picture(path, l, t, width=w)
    return slide.shapes.add_picture(path, l, t, height=h)


def metric_card(slide, l, t, w, big, label, color=ACCENT):
    card = slide.shapes.add_shape(1, l, t, w, Inches(1.6))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = color
    card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = big; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label; r2.font.size = Pt(13); r2.font.color.rgb = GREY


# =================================================================== 1
s = add_slide()
title_slide(s, 'Classificação de Sinais de Trânsito (GTSRB)',
            'Data Augmentation e Ensembles de Deep Learning')

# =================================================================== 2 agenda
s = content_slide('Agenda', kicker='Visão geral')
bullets(s, [
    '**O problema** — GTSRB, 43 classes, objetivo e estado da arte',
    '**Fase 1** — Pré-processamento e Data Augmentation',
    ('CLAHE e filtros clássicos · 7 estratégias · 2 arquiteturas', 1),
    ('Resultados e ablações (resolução, sampler)', 1),
    '**Fase 2** — Ensembles das redes treinadas',
    ('Métodos de agregação · saturação · matriz de confusão', 1),
    '**Conclusões** e comparação com o SOTA',
])

# =================================================================== 3 problema
s = content_slide('O problema', kicker='Contexto')
bullets(s, [
    '**GTSRB**: German Traffic Sign Recognition Benchmark',
    ('>50 000 imagens reais, **43 classes** de sinais', 1),
    ('Resoluções variáveis (~15–250 px), iluminação difícil, oclusões', 1),
    '**Objetivo**: maximizar a accuracy no conjunto de teste (12 630 imagens)',
    '**Estado da arte**: 99,82% (Cireşan et al., 2012 — ensemble de 25 redes)',
    ('Desempenho humano reportado: ~98,8%', 1),
], t=Inches(1.7), h=Inches(3.0))
metric_card(s, Inches(0.9),  Inches(5.3), Inches(3.4), '43',      'classes')
metric_card(s, Inches(4.85), Inches(5.3), Inches(3.4), '12 630',  'imagens de teste')
metric_card(s, Inches(8.8),  Inches(5.3), Inches(3.4), '99,82%',  'melhor resultado publicado')

# =================================================================== 4 CLAHE
s = content_slide('Pré-processamento — CLAHE e filtros clássicos', kicker='Fase 1')
bullets(s, [
    'Comparámos filtros clássicos: equalização global, sharpen, blur, **CLAHE**',
    '**CLAHE** (Contrast Limited Adaptive Histogram Equalization):',
    ('equaliza o histograma **localmente** (por blocos), com clip que evita amplificar ruído', 1),
    ('aplicado no canal de luminância (LAB), clip=2.0, tiles 4×4', 1),
    'É **determinístico** → aplicado de forma coerente em treino, validação e teste',
    '**Resultado**: técnicas clássicas de processamento de imagem ainda dão ganho real em deep learning',
])

# =================================================================== 5 estrategias
s = content_slide('7 estratégias de augmentation (complexidade crescente)', kicker='Fase 1')
rows = [
    ('Estratégia', 'Componentes acumulados', True),
    ('A_baseline',    'só Resize + Normalize (controlo)', False),
    ('B_geo_light',   '+ RandomAffine leve (rot. ±10°, transl.)', False),
    ('C_geo_heavy',   '+ scale, shear, RandomPerspective', False),
    ('D_photo',       '+ ColorJitter (brilho/contraste/sat./hue)', False),
    ('E_combo',       'C + D + RandomGaussianBlur', False),
    ('F_clahe_combo', 'E precedido de CLAHE', False),
    ('G_full',        'F + RandomErasing (cutout)', False),
]
tbl = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(1.55),
                         Inches(8.2), Inches(4.5)).table
tbl.columns[0].width = Inches(2.7); tbl.columns[1].width = Inches(5.5)
for ri, (a, b, hdr) in enumerate(rows):
    for ci, val in enumerate((a, b)):
        c = tbl.cell(ri, ci); c.text = val
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(13)
        pr.runs[0].font.bold = hdr or ci == 0
        pr.runs[0].font.color.rgb = WHITE if hdr else DARK
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT if hdr else (LIGHT if ri % 2 else WHITE)
txt(s, Inches(9.2), Inches(1.7), Inches(3.6), Inches(4.0),
    'Desenho incremental: cada estratégia acrescenta uma técnica, '
    'permitindo isolar o efeito de cada componente.\n\n'
    'Sem HorizontalFlip — inverteria sinais direcionais.\n\n'
    'Rotações ≤ ±15°.',
    size=15, color=GREY)

# =================================================================== 6 arquiteturas
s = content_slide('Arquiteturas e treino', kicker='Fase 1')
bullets(s, [
    '**Conv_VCPI** — CNN própria, estilo VGG (~332 k parâmetros)',
    ('3 blocos conv 3×3 + BN + ReLU + MaxPool + Dropout → GAP → FC', 1),
    '**ResNet18** — pré-treinada no ImageNet, transfer learning (~11 M parâmetros)',
], t=Inches(1.6), w=Inches(12), h=Inches(1.9), size=18, gap=8)
txt(s, Inches(0.7), Inches(3.35), Inches(6), Inches(0.5),
    'Receita de treino — 4 ingredientes:', size=16, bold=True, color=ACCENT)
bullets(s, [
    '**AdamW** — otimizador que ajusta os pesos para reduzir o erro',
    '**Redução do learning rate** — passo menor quando a validação estagna',
    '**Early stopping** — para quando a validação deixa de melhorar (evita overfitting)',
    '**Melhor checkpoint** — guarda a rede com menor erro de validação, não a última época',
], t=Inches(3.85), w=Inches(12.2), size=16, gap=8)

# =================================================================== 7 resultados f1
s = content_slide('Resultados da Fase 1', kicker='Fase 1')
image(s, f'{FIG}/phase1_bars.png', Inches(0.5), Inches(1.6), w=Inches(7.5))
bullets(s, [
    '**CLAHE dá ganho real**: E→F = +0,26 pp',
    'Geo + foto + CLAHE > qualquer componente isolado',
    'RandomErasing (G) não melhorou — regularização a mais',
    'A 48×48, a **Conv_VCPI bate a ResNet18** em todas as 7 estratégias',
    '**Melhor a 48×48**: Conv_VCPI + F = **99,48%**',
], l=Inches(8.2), t=Inches(1.7), w=Inches(4.8), size=15, gap=12)

# =================================================================== 8 ablacao resolucao
s = content_slide('Ablação decisiva — resolução de entrada', kicker='Fase 1')
bullets(s, [
    '**ResNet18 a 48×48**: ficava atrás da Conv_VCPI (98,98%)',
    '**ResNet18 a 224×224** (resolução nativa do ImageNet): **99,52%**',
    ('+0,54 pp vs 48×48 → passa a melhor modelo individual da Fase 1', 1),
    'Os pesos pré-treinados só rendem perto da escala em que foram aprendidos',
    ('transfer learning não é gratuito: depende do match de resolução/domínio', 1),
    'Conv_VCPI a 64×64: −0,10 pp ao dobro do tempo (50% das imagens < 48 px)',
    'WeightedRandomSampler: não melhorou — desbalanceamento não é o gargalo',
], t=Inches(1.7))

# =================================================================== 9 conclusoes f1
s = content_slide('Conclusões da Fase 1', kicker='Fase 1')
bullets(s, [
    '**99,52%** com um único modelo (ResNet18 + CLAHE a 224×224)',
    'Fatores de maior impacto: **CLAHE**, augmentation rico, match de resolução',
    'Erros remanescentes ~ todos por **ambiguidade visual**',
    ('sinais triangulares parecidos, limites de velocidade adjacentes, imagens degradadas', 1),
    'Track leakage detetado: validação satura a ~100% (não discrimina modelos)',
    '→ Perto do tecto de um único modelo. **Próximo passo: ensembles.**',
])

# =================================================================== 10 ensemble ideia
s = content_slide('Fase 2 — porquê ensembles?', kicker='Fase 2')
bullets(s, [
    'Modelos **diversos** (arquiteturas e augmentations diferentes) erram em imagens diferentes',
    'Combinar as previsões cancela erros independentes → melhora a accuracy',
    '**Pool de 17 modelos** treinados na Fase 1:',
    ('14 principais (7 estratégias × 2 arquiteturas)', 1),
    ('+ 3 ablações (Conv@64, Conv+WeightedSampler, ResNet18@224)', 1),
    'Logits calculados uma vez e guardados em cache → avaliar combinações é instantâneo',
], t=Inches(1.7))

# =================================================================== 11 metodos
s = content_slide('Métodos de agregação', kicker='Fase 2')
bullets(s, [
    '**Soft voting** — média das probabilidades softmax e argmax (preserva a confiança)',
    '**Weighted** — média ponderada pela accuracy individual de cada modelo',
    '**Hard voting** — cada modelo vota na sua classe argmax; ganha a maioria',
], t=Inches(1.8), size=19, gap=16)
txt(s, Inches(0.7), Inches(4.2), Inches(12), Inches(1),
    'Como os logits estão em cache, comparar métodos é só fazer contas sobre tensores.',
    size=15, color=GREY)

# =================================================================== 12 resultados f2
s = content_slide('Resultados — o ensemble bate o melhor modelo', kicker='Fase 2')
image(s, f'{FIG}/phase2_methods.png', Inches(0.5), Inches(1.7), w=Inches(7.4))
bullets(s, [
    'Melhor individual: **99,52%** (61 erros)',
    'Hard voting: 99,76%',
    'Weighted: 99,79%',
    '**Soft voting (17): 99,80%** — 25 erros',
    'Soft > weighted > hard',
    ('hard perde por descartar a confiança', 1),
], l=Inches(8.1), t=Inches(1.8), w=Inches(4.9), size=16, gap=10)

# =================================================================== 13 saturacao
s = content_slide('Quantos modelos? Diversidade > quantidade', kicker='Fase 2')
image(s, f'{FIG}/phase2_saturation.png', Inches(0.5), Inches(1.7), w=Inches(7.6))
bullets(s, [
    'Modelos adicionados por ordem de accuracy individual (sem cherry-picking)',
    'A accuracy **satura a partir de ~13 modelos**',
    'O ganho vem da **diversidade**, não do número de redes',
], l=Inches(8.2), t=Inches(2.0), w=Inches(4.8), size=16, gap=12)

# =================================================================== 14 confusao
s = content_slide('Matriz de confusão do ensemble', kicker='Fase 2')
image(s, f'{FIG}/phase2_confusion.png', Inches(0.5), Inches(1.7), w=Inches(12.3))

# =================================================================== 15 conclusao
s = content_slide('Conclusão global', kicker='Síntese')
bullets(s, [
    '**Processamento de imagem importa**: CLAHE deu ganho real; melhor estratégia = geo + foto + CLAHE',
    '**Resolução é decisiva no transfer learning**: ResNet18 só ganha à resolução nativa (224×224)',
    '**Fase 1** — melhor modelo individual: **99,52%**',
    '**Fase 2** — ensemble soft voting de redes diversas: **99,80%** (+0,28 pp)',
    ('ao nível do SOTA (99,82%) — ~2 imagens de diferença; o SOTA é também um ensemble', 1),
    'O ganho vem da **diversidade** das redes e **satura a poucos modelos** (~13)',
    '**Trabalho futuro**: Test-Time Augmentation · split track-aware · Spatial Transformer Networks',
], t=Inches(1.55), w=Inches(12.4), size=17, gap=11)

# =================================================================== 16 fim
s = add_slide()
title_slide(s, 'Obrigado!', 'Perguntas?')

out = 'apresentacao_VCPI.pptx'
prs.save(out)
print('guardado:', out, '| slides:', len(prs.slides._sldIdLst))
